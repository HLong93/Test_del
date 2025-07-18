#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
知识库文档预处理脚本
支持处理Word、PDF、Excel、PowerPoint文档
提取文本和图片内容，转换为Q&A格式
"""

import os
import sys
import re
import json
from pathlib import Path
from typing import List, Dict, Any, Tuple
import logging

# 文档处理库
try:
    import docx
    from docx.document import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
except ImportError:
    print("请安装python-docx: pip install python-docx")
    sys.exit(1)

try:
    import PyPDF2
    import fitz  # PyMuPDF
except ImportError:
    print("请安装PDF处理库: pip install PyPDF2 PyMuPDF")
    sys.exit(1)

try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    print("请安装openpyxl: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("请安装python-pptx: pip install python-pptx")
    sys.exit(1)

# 图像处理
try:
    from PIL import Image
    import base64
    from io import BytesIO
except ImportError:
    print("请安装Pillow: pip install Pillow")
    sys.exit(1)

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """文档处理器主类"""
    
    def __init__(self, input_dir: str, output_dir: str):
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
        # 支持的文件格式
        self.supported_formats = {
            '.docx': self.process_word,
            '.doc': self.process_word,
            '.pdf': self.process_pdf,
            '.xlsx': self.process_excel,
            '.xls': self.process_excel,
            '.pptx': self.process_powerpoint,
            '.ppt': self.process_powerpoint
        }
        
        self.processed_count = 0
        
    def process_all_documents(self):
        """处理所有文档"""
        logger.info(f"开始处理目录: {self.input_dir}")
        
        for root, dirs, files in os.walk(self.input_dir):
            for file in files:
                file_path = Path(root) / file
                if file_path.suffix.lower() in self.supported_formats:
                    try:
                        logger.info(f"处理文件: {file_path}")
                        self.process_single_document(file_path)
                        self.processed_count += 1
                    except Exception as e:
                        logger.error(f"处理文件 {file_path} 时出错: {str(e)}")
                        
        logger.info(f"处理完成，共处理 {self.processed_count} 个文件")
        
    def process_single_document(self, file_path: Path):
        """处理单个文档"""
        file_ext = file_path.suffix.lower()
        processor = self.supported_formats.get(file_ext)
        
        if processor:
            content = processor(file_path)
            if content:
                self.save_processed_content(file_path, content)
        else:
            logger.warning(f"不支持的文件格式: {file_ext}")
            
    def process_word(self, file_path: Path) -> Dict[str, Any]:
        """处理Word文档"""
        try:
            doc = docx.Document(file_path)
            content = {
                'title': file_path.stem,
                'source_file': str(file_path),
                'type': 'word',
                'sections': [],
                'images': []
            }
            
            current_section = {
                'heading': '',
                'content': '',
                'tables': [],
                'images': []
            }
            
            # 处理段落和表格
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    paragraph = Paragraph(element, doc)
                    text = paragraph.text.strip()
                    
                    if text:
                        # 检查是否为标题
                        if paragraph.style.name.startswith('Heading') or self.is_heading(text):
                            if current_section['content'] or current_section['tables']:
                                content['sections'].append(current_section.copy())
                            current_section = {
                                'heading': text,
                                'content': '',
                                'tables': [],
                                'images': []
                            }
                        else:
                            current_section['content'] += text + '\n'
                            
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    table_data = self.extract_table_data(table)
                    current_section['tables'].append(table_data)
            
            # 添加最后一个section
            if current_section['content'] or current_section['tables']:
                content['sections'].append(current_section)
                
            # 提取图片
            content['images'] = self.extract_word_images(doc)
            
            return content
            
        except Exception as e:
            logger.error(f"处理Word文档 {file_path} 时出错: {str(e)}")
            return None
            
    def process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """处理PDF文档"""
        try:
            content = {
                'title': file_path.stem,
                'source_file': str(file_path),
                'type': 'pdf',
                'sections': [],
                'images': []
            }
            
            # 使用PyMuPDF处理PDF
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取文本
                text = page.get_text()
                if text.strip():
                    section = {
                        'heading': f'第{page_num + 1}页',
                        'content': text,
                        'tables': [],
                        'images': []
                    }
                    content['sections'].append(section)
                
                # 提取图片
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            img_info = {
                                'page': page_num + 1,
                                'index': img_index,
                                'data': base64.b64encode(img_data).decode(),
                                'format': 'png'
                            }
                            content['images'].append(img_info)
                        pix = None
                    except Exception as e:
                        logger.warning(f"提取PDF图片时出错: {str(e)}")
            
            doc.close()
            return content
            
        except Exception as e:
            logger.error(f"处理PDF文档 {file_path} 时出错: {str(e)}")
            return None
            
    def process_excel(self, file_path: Path) -> Dict[str, Any]:
        """处理Excel文档"""
        try:
            workbook = load_workbook(file_path, data_only=True)
            content = {
                'title': file_path.stem,
                'source_file': str(file_path),
                'type': 'excel',
                'sections': [],
                'images': []
            }
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # 提取工作表数据
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_data.append([str(cell) if cell is not None else '' for cell in row])
                
                if sheet_data:
                    section = {
                        'heading': f'工作表: {sheet_name}',
                        'content': '',
                        'tables': [{'data': sheet_data}],
                        'images': []
                    }
                    content['sections'].append(section)
            
            return content
            
        except Exception as e:
            logger.error(f"处理Excel文档 {file_path} 时出错: {str(e)}")
            return None
            
    def process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """处理PowerPoint文档"""
        try:
            prs = Presentation(file_path)
            content = {
                'title': file_path.stem,
                'source_file': str(file_path),
                'type': 'powerpoint',
                'sections': [],
                'images': []
            }
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = {
                    'heading': f'幻灯片 {slide_num + 1}',
                    'content': '',
                    'tables': [],
                    'images': []
                }
                
                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content['content'] += shape.text + '\n'
                    
                    # 提取图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            img_data = base64.b64encode(image.blob).decode()
                            img_info = {
                                'slide': slide_num + 1,
                                'data': img_data,
                                'format': image.ext
                            }
                            content['images'].append(img_info)
                        except Exception as e:
                            logger.warning(f"提取PPT图片时出错: {str(e)}")
                
                if slide_content['content'].strip():
                    content['sections'].append(slide_content)
            
            return content
            
        except Exception as e:
            logger.error(f"处理PowerPoint文档 {file_path} 时出错: {str(e)}")
            return None

    def is_heading(self, text: str) -> bool:
        """判断文本是否为标题"""
        # 简单的标题判断逻辑
        heading_patterns = [
            r'^第[一二三四五六七八九十\d]+[章节部分]',
            r'^\d+[\.\s]',
            r'^[一二三四五六七八九十]+[\.\s]',
            r'^问题[:：]',
            r'^解决方案[:：]',
            r'^步骤[:：]',
            r'^注意[:：]'
        ]

        for pattern in heading_patterns:
            if re.match(pattern, text.strip()):
                return True

        # 如果文本很短且包含关键词，可能是标题
        if len(text.strip()) < 50 and any(keyword in text for keyword in ['问题', '解决', '步骤', '方法', '设置', '配置']):
            return True

        return False

    def extract_table_data(self, table: Table) -> Dict[str, Any]:
        """提取表格数据"""
        table_data = {
            'headers': [],
            'rows': []
        }

        for i, row in enumerate(table.rows):
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())

            if i == 0:
                table_data['headers'] = row_data
            else:
                table_data['rows'].append(row_data)

        return table_data

    def extract_word_images(self, doc: Document) -> List[Dict[str, Any]]:
        """从Word文档中提取图片"""
        images = []

        try:
            # 获取文档中的所有图片关系
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        img_info = {
                            'data': base64.b64encode(image_data).decode(),
                            'format': rel.target_ref.split('.')[-1],
                            'relation_id': rel.rId
                        }
                        images.append(img_info)
                    except Exception as e:
                        logger.warning(f"提取Word图片时出错: {str(e)}")
        except Exception as e:
            logger.warning(f"处理Word图片关系时出错: {str(e)}")

        return images

    def convert_to_qa_format(self, content: Dict[str, Any]) -> List[Dict[str, Any]]:
        """将内容转换为Q&A格式"""
        qa_pairs = []

        # 基于文档标题生成基础上下文
        base_context = f"文档来源：{content['title']}"

        for section in content['sections']:
            heading = section['heading']
            text_content = section['content']
            tables = section.get('tables', [])

            # 处理文本内容
            if text_content.strip():
                # 尝试识别问题和答案
                qa_pair = self.extract_qa_from_text(heading, text_content, base_context)
                if qa_pair:
                    qa_pairs.append(qa_pair)

            # 处理表格内容
            for table in tables:
                table_qa = self.extract_qa_from_table(heading, table, base_context)
                if table_qa:
                    qa_pairs.append(table_qa)

        # 处理图片内容
        for i, image in enumerate(content.get('images', [])):
            image_qa = self.create_image_qa(content['title'], image, i, base_context)
            if image_qa:
                qa_pairs.append(image_qa)

        return qa_pairs

    def extract_qa_from_text(self, heading: str, content: str, base_context: str) -> Dict[str, Any]:
        """从文本中提取Q&A"""
        # 清理内容
        content = re.sub(r'\n+', '\n', content.strip())

        # 尝试识别问题模式
        question_patterns = [
            r'问题[:：](.+?)(?=解决|答案|方法|步骤|$)',
            r'故障[:：](.+?)(?=解决|修复|处理|$)',
            r'错误[:：](.+?)(?=解决|修复|处理|$)',
            r'如何(.+?)(?=\n|$)',
            r'怎样(.+?)(?=\n|$)'
        ]

        question = heading if heading else "相关问题"
        answer = content

        # 尝试从内容中提取更具体的问题
        for pattern in question_patterns:
            match = re.search(pattern, content, re.DOTALL | re.IGNORECASE)
            if match:
                question = match.group(1).strip()
                break

        # 生成通用化的问题
        generalized_question = self.generalize_question(question, content)

        return {
            'question': generalized_question,
            'answer': self.enhance_answer(answer, base_context),
            'keywords': self.extract_keywords(content),
            'source': base_context,
            'type': 'text'
        }

    def extract_qa_from_table(self, heading: str, table: Dict[str, Any], base_context: str) -> Dict[str, Any]:
        """从表格中提取Q&A"""
        headers = table.get('headers', [])
        rows = table.get('rows', [])

        if not headers or not rows:
            return None

        # 构建表格内容描述
        table_content = f"表格标题：{heading}\n"
        table_content += "表格内容：\n"

        # 添加表头
        table_content += " | ".join(headers) + "\n"
        table_content += "|".join(["---"] * len(headers)) + "\n"

        # 添加数据行
        for row in rows:
            table_content += " | ".join(row) + "\n"

        question = f"关于{heading}的详细信息"

        return {
            'question': question,
            'answer': table_content,
            'keywords': headers + [heading],
            'source': base_context,
            'type': 'table'
        }

    def create_image_qa(self, doc_title: str, image: Dict[str, Any], index: int, base_context: str) -> Dict[str, Any]:
        """为图片创建Q&A"""
        question = f"{doc_title}中的图片{index + 1}"

        # 简单的图片描述（实际应用中可以集成图像识别API）
        answer = f"这是来自文档《{doc_title}》的第{index + 1}张图片。"

        # 根据图片格式添加描述
        img_format = image.get('format', 'unknown')
        if 'page' in image:
            answer += f"位于第{image['page']}页。"
        elif 'slide' in image:
            answer += f"位于第{image['slide']}张幻灯片。"

        return {
            'question': question,
            'answer': answer,
            'keywords': [doc_title, '图片', '图像'],
            'source': base_context,
            'type': 'image',
            'image_data': image['data'],
            'image_format': img_format
        }

    def generalize_question(self, question: str, content: str) -> str:
        """将具体问题泛化为通用问题"""
        # 移除具体的数字、日期、名称等
        generalized = question

        # 替换具体数字为通用描述
        generalized = re.sub(r'\d{4}-\d{2}-\d{2}', '[日期]', generalized)
        generalized = re.sub(r'\d+', '[数字]', generalized)

        # 替换具体名称为通用描述
        generalized = re.sub(r'小区\d+', '[小区名称]', generalized)
        generalized = re.sub(r'用户\w+', '[用户名]', generalized)

        # 添加常见问题前缀
        if not any(prefix in generalized for prefix in ['如何', '怎样', '什么', '为什么', '问题']):
            if '设置' in content or '配置' in content:
                generalized = f"如何{generalized}"
            elif '错误' in content or '故障' in content:
                generalized = f"如何解决{generalized}"
            elif '步骤' in content or '流程' in content:
                generalized = f"如何进行{generalized}"

        return generalized

    def enhance_answer(self, answer: str, base_context: str) -> str:
        """增强答案内容"""
        enhanced = f"{base_context}\n\n{answer}"

        # 确保SQL语句被保留
        sql_pattern = r'(SELECT|INSERT|UPDATE|DELETE|CREATE|ALTER|DROP)\s+.*?;'
        sql_matches = re.findall(sql_pattern, answer, re.IGNORECASE | re.DOTALL)

        if sql_matches:
            enhanced += "\n\n相关SQL语句：\n"
            for sql in sql_matches:
                enhanced += f"```sql\n{sql}\n```\n"

        # 确保操作步骤被突出显示
        if '步骤' in answer:
            steps = re.findall(r'步骤\s*\d+[:：](.+?)(?=步骤\s*\d+|$)', answer, re.DOTALL)
            if steps:
                enhanced += "\n\n操作步骤：\n"
                for i, step in enumerate(steps, 1):
                    enhanced += f"{i}. {step.strip()}\n"

        return enhanced

    def extract_keywords(self, content: str) -> List[str]:
        """提取关键词"""
        # 技术关键词
        tech_keywords = [
            'SRMS', 'SQL', 'IE', 'Edge', 'Login', 'Batch', 'Journal',
            'Demand Note', 'SPS', 'Occupant', 'Building ID', 'CashType',
            'Synergis', 'Community App', 'Facility Booking'
        ]

        # 操作关键词
        action_keywords = [
            '设置', '配置', '登录', '删除', '更新', '导出', '列印', '解锁',
            '修复', '处理', '解决', '安装', '注册', '上传', '下载'
        ]

        found_keywords = []
        content_lower = content.lower()

        for keyword in tech_keywords + action_keywords:
            if keyword.lower() in content_lower:
                found_keywords.append(keyword)

        # 提取其他重要词汇
        important_words = re.findall(r'[\u4e00-\u9fff]{2,}', content)
        found_keywords.extend(list(set(important_words))[:10])  # 限制数量

        return list(set(found_keywords))

    def save_processed_content(self, original_file: Path, content: Dict[str, Any]):
        """保存处理后的内容"""
        # 转换为Q&A格式
        qa_pairs = self.convert_to_qa_format(content)

        if not qa_pairs:
            logger.warning(f"文件 {original_file} 没有生成有效的Q&A内容")
            return

        # 生成输出文件名
        output_filename = f"{self.processed_count:02d}_{original_file.stem}_processed.docx"
        output_path = self.output_dir / output_filename

        # 创建Word文档
        self.create_word_document(qa_pairs, output_path, content['title'])

        # 保存JSON格式的原始数据（用于调试）
        json_path = self.output_dir / f"{self.processed_count:02d}_{original_file.stem}_data.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump({
                'original_file': str(original_file),
                'qa_pairs': qa_pairs,
                'raw_content': content
            }, f, ensure_ascii=False, indent=2)

        logger.info(f"已保存处理结果: {output_path}")

    def create_word_document(self, qa_pairs: List[Dict[str, Any]], output_path: Path, title: str):
        """创建Word文档"""
        doc = docx.Document()

        # 添加标题
        doc.add_heading(f'知识库文档：{title}', 0)

        # 添加说明
        doc.add_paragraph(f'本文档包含 {len(qa_pairs)} 个问答对，来源于原始文档的处理和整理。')

        for i, qa in enumerate(qa_pairs, 1):
            # 添加问题
            question_para = doc.add_heading(f'Q{i}: {qa["question"]}', level=2)

            # 添加答案
            answer_para = doc.add_paragraph()
            answer_para.add_run('A: ').bold = True
            answer_para.add_run(qa['answer'])

            # 如果有图片，添加图片
            if qa.get('type') == 'image' and qa.get('image_data'):
                try:
                    # 解码base64图片数据
                    img_data = base64.b64decode(qa['image_data'])
                    img_stream = BytesIO(img_data)

                    # 添加图片到文档
                    doc.add_picture(img_stream, width=docx.shared.Inches(4))
                except Exception as e:
                    logger.warning(f"添加图片到Word文档时出错: {str(e)}")

            # 添加关键词
            if qa.get('keywords'):
                keywords_para = doc.add_paragraph()
                keywords_para.add_run('关键词: ').italic = True
                keywords_para.add_run(', '.join(qa['keywords']))

            # 添加分隔线
            doc.add_paragraph('─' * 50)

        # 保存文档
        doc.save(output_path)


def main():
    """主函数"""
    input_directory = "待处理知识库"
    output_directory = "已处理知识库"

    if not os.path.exists(input_directory):
        logger.error(f"输入目录不存在: {input_directory}")
        return

    processor = DocumentProcessor(input_directory, output_directory)
    processor.process_all_documents()

    print(f"\n处理完成！")
    print(f"输入目录: {input_directory}")
    print(f"输出目录: {output_directory}")
    print(f"处理文件数: {processor.processed_count}")


if __name__ == "__main__":
    main()
